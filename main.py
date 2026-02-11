import os, subprocess, time
from rich.console import Console
from rich.markdown import Markdown
from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn, TimeElapsedColumn
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

DATA_DIR = "data"
MODEL_NAME = "llama3.1"
BOT_LANGUAGE = "Italian"
BOT_NAME = "Ares"
PROMPT_FILE = "prompt.txt"
console = Console()

# -------------------- Lettori per tipologia --------------------
def read_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def read_pdf(path):
    reader = PdfReader(path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)

def read_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# -------------------- Funzione principale di caricamento con barra di avanzamento --------------------
def load_data():
    # Raccogli prima tutti i file validi da processare
    files_to_process = []
    for filename in os.listdir(DATA_DIR):
        path = os.path.join(DATA_DIR, filename)
        if os.path.isfile(path):
            ext = filename.lower().split('.')[-1]
            if ext in ["txt", "docx", "pdf", "pptx"]:
                files_to_process.append((filename, path, ext))
    
    if not files_to_process:
        console.print("[bold yellow]⚠️ Nessun file valido trovato nella cartella 'data'.[/bold yellow]")
        return ""
    
    texts = []
    
    # Barra di avanzamento avanzata con tempo e file corrente
    with Progress(
        SpinnerColumn(style="cyan"),
        TextColumn("[bold blue]{task.description}"),
        BarColumn(bar_width=40, complete_style="green"),
        TextColumn("[bold]{task.completed}/{task.total}"),
        TimeElapsedColumn(),
        TextColumn("[cyan]{task.fields[filename]}"),
        console=console,
        transient=False
    ) as progress:
        task = progress.add_task("Caricamento documenti...", total=len(files_to_process), filename="Inizializzazione...")
        
        for filename, path, ext in files_to_process:
            progress.update(task, filename=f"📄 {filename}")
            try:
                if ext == "txt":
                    content = read_txt(path)
                elif ext == "docx":
                    content = read_docx(path)
                elif ext == "pdf":
                    content = read_pdf(path)
                elif ext == "pptx":
                    content = read_pptx(path)
                else:
                    try:
                        content = read_txt(path)
                    except:
                        console.print(f"[bold red]❌ File non supportato: {filename}[/bold red]")
                        progress.update(task, advance=1)
                        continue
                
                texts.append(f"=== {filename} ===\n{content}")
                progress.update(task, advance=1)
                
            except Exception as e:
                console.print(f"[bold red]❌ Errore leggendo {filename}: {e}[/bold red]")
                progress.update(task, advance=1)
                continue
    
    console.print(f"\n[bold green]✅ Caricamento completato: {len(texts)}/{len(files_to_process)} file elaborati[/bold green]")
    return "\n\n".join(texts)

# -------------------- Interrogazione Ollama --------------------
def ask_ollama(context, question):
    prompt = open(PROMPT_FILE, "r").read().format(
        BOT_NAME=BOT_NAME,
        BOT_LANGUAGE=BOT_LANGUAGE,
        context=context,
        question=question
    )

    result = subprocess.run(
        ["ollama", "run", MODEL_NAME],
        input=prompt,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="ignore"
    )
    return result.stdout.strip()

# -------------------- Main --------------------
def main():
    console.print("\n[bold cyan]📂 Caricamento dati in corso...[/bold cyan]\n")
    context = load_data()
    
    if not context.strip():
        return

    console.print("\n[bold green]✨ Pronto per le tue domande![/bold green]\n")

    while True:
        question = input("❓ Domanda (scrivi 'exit' per uscire): ")
        if question.lower() == "exit":
            break

        console.print("\n[bold cyan]⏳ Elaborazione in corso...[/bold cyan]")
        start_time = time.time()
        answer = ask_ollama(context, question)
        elapsed = time.time() - start_time
        
        console.print(f"\n🤖 Risposta [dim](elaborata in {elapsed:.2f}s)[/dim]:")
        console.print(Markdown(answer))
        console.print()

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        console.print("\n[bold yellow]⚠️ Interrotto dall'utente[/bold yellow]")
        pass